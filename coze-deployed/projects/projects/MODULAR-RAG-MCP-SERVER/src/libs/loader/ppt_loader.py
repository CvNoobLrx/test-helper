"""PPT Loader implementation using MarkItDown and python-pptx.

This module implements PPTX parsing with image extraction support,
converting PowerPoint slides to standardized Markdown format with image placeholders.

Features:
- Text extraction and Markdown conversion via MarkItDown
- Embedded image extraction per slide via python-pptx
- Image placeholder insertion with metadata tracking
- Graceful degradation if image extraction fails
"""

from __future__ import annotations

import hashlib
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from markitdown import MarkItDown
    MARKITDOWN_AVAILABLE = True
except ImportError:
    MARKITDOWN_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from PIL import Image
import io

from src.core.types import Document
from src.libs.loader.base_loader import BaseLoader

logger = logging.getLogger(__name__)


class PptLoader(BaseLoader):
    """PPT Loader using MarkItDown for text extraction and python-pptx for images.

    This loader:
    1. Extracts text from PPTX and converts to Markdown
    2. Extracts embedded images and saves to data/images/{doc_hash}/
    3. Inserts image placeholders in the format [IMAGE: {image_id}]
    4. Records image metadata in Document.metadata.images

    Graceful Degradation:
        If image extraction fails, logs warning and continues with text-only parsing.
    """

    SUPPORTED_EXTENSIONS = {".pptx"}

    def __init__(
        self,
        extract_images: bool = True,
        image_storage_dir: str | Path = "data/images",
    ):
        if not MARKITDOWN_AVAILABLE:
            raise ImportError(
                "MarkItDown is required for PptLoader. "
                "Install with: pip install markitdown"
            )
        self.extract_images = extract_images
        self.image_storage_dir = Path(image_storage_dir)
        self._markitdown = MarkItDown()

    def load(self, file_path: str | Path) -> Document:
        path = self._validate_file(file_path)
        if path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"File is not a PPTX: {path}")

        doc_hash = self._compute_file_hash(path)
        doc_id = f"doc_{doc_hash[:16]}"

        # Parse PPTX with MarkItDown
        try:
            result = self._markitdown.convert(str(path))
            text_content = result.text_content if hasattr(result, "text_content") else str(result)
        except Exception as e:
            logger.error(f"Failed to parse PPTX {path}: {e}")
            raise RuntimeError(f"PPTX parsing failed: {e}") from e

        metadata: Dict[str, Any] = {
            "source_path": str(path),
            "doc_type": "pptx",
            "doc_hash": doc_hash,
        }

        title = self._extract_title(text_content)
        if title:
            metadata["title"] = title

        # Handle image extraction
        if self.extract_images:
            try:
                text_content, images_metadata = self._extract_and_process_images(
                    path, text_content, doc_hash
                )
                if images_metadata:
                    metadata["images"] = images_metadata
            except Exception as e:
                logger.warning(
                    f"Image extraction failed for {path}, continuing with text-only: {e}"
                )

        return Document(id=doc_id, text=text_content, metadata=metadata)

    def _compute_file_hash(self, file_path: Path) -> str:
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    def _extract_title(self, text: str) -> Optional[str]:
        lines = text.split("\n")
        for line in lines[:20]:
            line = line.strip()
            if line.startswith("# "):
                return line[2:].strip()
        for line in lines[:10]:
            line = line.strip()
            if line:
                return line
        return None

    def _extract_and_process_images(
        self,
        pptx_path: Path,
        text_content: str,
        doc_hash: str,
    ) -> tuple[str, List[Dict[str, Any]]]:
        if not self.extract_images:
            return text_content, []

        if not PPTX_AVAILABLE:
            logger.warning(f"python-pptx not available, skipping image extraction for {pptx_path}")
            return text_content, []

        images_metadata: List[Dict[str, Any]] = []
        modified_text = text_content

        try:
            image_dir = self.image_storage_dir / doc_hash
            image_dir.mkdir(parents=True, exist_ok=True)

            prs = Presentation(str(pptx_path))

            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if not shape.has_table and hasattr(shape, "image"):
                        try:
                            image_blob = shape.image.blob
                            content_type = shape.image.content_type
                            ext = content_type.split("/")[-1] if "/" in content_type else "png"
                            if ext == "jpeg":
                                ext = "jpg"

                            image_id = self._generate_image_id(doc_hash, slide_idx + 1, shape_idx + 1)
                            image_filename = f"{image_id}.{ext}"
                            image_path = image_dir / image_filename

                            with open(image_path, "wb") as img_file:
                                img_file.write(image_blob)

                            try:
                                img = Image.open(io.BytesIO(image_blob))
                                width, height = img.size
                            except Exception:
                                width, height = 0, 0

                            placeholder = f"[IMAGE: {image_id}]"
                            insert_position = len(modified_text)
                            modified_text += f"\n{placeholder}\n"

                            try:
                                relative_path = image_path.relative_to(Path.cwd())
                            except ValueError:
                                relative_path = image_path.absolute()

                            image_metadata = {
                                "id": image_id,
                                "path": str(relative_path),
                                "page": slide_idx + 1,
                                "text_offset": insert_position + 1,
                                "text_length": len(placeholder),
                                "position": {
                                    "width": width,
                                    "height": height,
                                    "page": slide_idx + 1,
                                    "index": shape_idx,
                                },
                            }
                            images_metadata.append(image_metadata)
                            logger.debug(f"Extracted image {image_id} from slide {slide_idx + 1}")

                        except Exception as e:
                            logger.warning(
                                f"Failed to extract image {shape_idx} from slide {slide_idx + 1}: {e}"
                            )
                            continue

            if images_metadata:
                logger.info(f"Extracted {len(images_metadata)} images from {pptx_path}")

            return modified_text, images_metadata

        except Exception as e:
            logger.warning(f"Image extraction failed for {pptx_path}: {e}")
            return text_content, []

    @staticmethod
    def _generate_image_id(doc_hash: str, slide: int, sequence: int) -> str:
        return f"{doc_hash[:8]}_{slide}_{sequence}"
